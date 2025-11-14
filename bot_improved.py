import fitz  # PyMuPDF
from pptx import Presentation
from telegram import Update
from telegram.ext import ApplicationBuilder, MessageHandler, filters, ContextTypes
import os
from flask import Flask
import threading
import asyncio
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer

# --- جلب التوكن من Environment أو مباشرة ---
BOT_TOKEN = "7935681061:AAG6zPjZ_0mifx_Mccijvjzzu_cFVFWrKaw"

# --- إعداد Web Server ---
app_web = Flask(__name__)

@app_web.route("/")
def home():
    return "Bot is running on Render!"

def run_web():
    port = int(os.environ.get("PORT", 10000))
    app_web.run(host="0.0.0.0", port=port)

# --- دالة تلخيص النص ---
def summarize_text_local(text, sentences_count=2):
    if not text.strip():
        return "لا يوجد نص للسلايد"
    
    parser = PlaintextParser.from_string(text, Tokenizer("arabic"))
    summarizer = LsaSummarizer()
    summary_sentences = summarizer(parser.document, sentences_count)
    summary = " ".join([str(sentence) for sentence in summary_sentences])
    return summary if summary else "لا يوجد محتوى كافي للتلخيص"

# --- معالجة PPT ---
def process_ppt(file_path):
    prs = Presentation(file_path)
    summaries = []
    for i, slide in enumerate(prs.slides, start=1):
        text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        summary = summarize_text_local(text)
        summaries.append(f"📌 سلايد {i}:\n{summary}\n")
    return summaries

# --- معالجة PDF ---
def process_pdf(file_path):
    doc = fitz.open(file_path)
    summaries = []
    for i, page in enumerate(doc, start=1):
        text = page.get_text()
        summary = summarize_text_local(text)
        summaries.append(f"📌 صفحة {i}:\n{summary}\n")
    return summaries

# --- استقبال الملفات ---
async def handle_file(update: Update, context: ContextTypes.DEFAULT_TYPE):
    document = update.message.document
    file_name = document.file_name
    file_path = f"temp_{file_name}"

    file = await document.get_file()
    await file.download_to_drive(file_path)

    try:
        if file_name.endswith(".pptx"):
            summaries = process_ppt(file_path)
        elif file_name.endswith(".pdf"):
            summaries = process_pdf(file_path)
        else:
            await update.message.reply_text("❌ نوع الملف غير مدعوم، استخدم PDF أو PPTX.")
            return
    except Exception as e:
        await update.message.reply_text(f"⚠️ حدث خطأ أثناء معالجة الملف: {e}")
        return
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)

    # إرسال كل الملخصات
    for summary in summaries:
        for i in range(0, len(summary), 4000):
            await update.message.reply_text(summary[i:i+4000])

    # إنشاء ملف TXT بالملخصات
    summary_file = f"summary_{file_name}.txt"
    with open(summary_file, "w", encoding="utf-8") as f:
        for s in summaries:
            f.write(s + "\n")

    await update.message.reply_document(open(summary_file, "rb"))
    os.remove(summary_file)

# --- تشغيل البوت ---
async def run_bot():
    application = ApplicationBuilder().token(BOT_TOKEN).build()
    application.add_handler(MessageHandler(filters.Document.ALL, handle_file))
    await application.initialize()   # تهيئة البوت
    await application.start()        # تشغيل البوت
    await application.updater.start_polling()  # بدء استقبال الرسائل
    await asyncio.Event().wait()     # إبقاء البوت يعمل

# --- Main ---
if __name__ == "__main__":
    # تشغيل Flask في Thread منفصل
    threading.Thread(target=run_web, daemon=True).start()

    # تشغيل Telegram bot بدون asyncio.run
    loop = asyncio.get_event_loop()
    loop.create_task(run_bot())
    loop.run_forever()
